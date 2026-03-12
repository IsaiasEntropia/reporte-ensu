# Librerías

from __future__ import annotations

import os
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Tuple, Optional

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Funciones

# ---------- Especificaciones de layout (posiciones fijas) ----------

@dataclass(frozen=True)
class TextSpec:
    col: str                  # columna en df_main
    left: int
    top: int
    width: int
    height: int
    font_size: int = 24
    bold: bool = False
    font_name: str = "Atkinson Hyperlegible"
    font_color: str = "2E4D58"
    fmt: str = "{val}"        


@dataclass(frozen=True)
class TopItemSpec:
    # Para los top-5 (longtable) pondremos imagen + texto por slot fijo
    img_left: int
    img_top: int
    img_width: int
    img_height: int
    txt_left: int
    txt_top: int
    txt_width: int
    txt_height: int
    font_size: int = 14
    bold: bool = False
    font_name: str = "Atkinson Hyperlegible"
    fmt: str = "{score}"   # label y score vienen del longtable



# Ejemplo: 3 textos fijos en un slide (ajusta posiciones)
FIXED_TEXT_FIELDS: List[TextSpec] = [
    TextSpec(col="p_inseg", 
             left=Inches(0.4),
             top=Inches(4.8), 
             width=Inches(1.3), 
             height=Inches(0.5), 
             font_size=26, 
             bold=True),
    TextSpec(col='Muy o algo efectivo', 
             left=Inches(4.84), 
             top=Inches(3.87), 
             width=Inches(0.85), 
             height=Inches(0.5), 
             font_size=19),
    TextSpec(col="Poco o nada efectivo", 
             left=Inches(4.84), 
             top=Inches(4.4), 
             width=Inches(0.85), 
             height=Inches(0.5), 
             font_size=19),
]


# Ejemplo: 5 slots fijos para top-5 (imagen + texto debajo)
TOP5_SLOTS: List[TopItemSpec] = [
    TopItemSpec(
        img_left=Inches(4.5), img_top=Inches(1.65), img_width=Inches(0.44), img_height=Inches(0.44),
        txt_left=Inches(4.48), txt_top=Inches(2.16), txt_width=Inches(0.65), txt_height=Inches(0.5),
        font_size=12, bold=True
    ),
    TopItemSpec(
        img_left=Inches(5.41), img_top=Inches(1.65), img_width=Inches(0.44), img_height=Inches(0.44),
        txt_left=Inches(5.43), txt_top=Inches(2.16), txt_width=Inches(0.65), txt_height=Inches(0.5),
        font_size=12, bold=True
    ),
    TopItemSpec(
        img_left=Inches(6.47), img_top=Inches(1.65), img_width=Inches(0.44), img_height=Inches(0.44),
        txt_left=Inches(6.43), txt_top=Inches(2.16), txt_width=Inches(0.65), txt_height=Inches(0.5),
        font_size=12, bold=True
    ),
    TopItemSpec(
        img_left=Inches(7.49), img_top=Inches(1.65), img_width=Inches(0.44), img_height=Inches(0.44),
        txt_left=Inches(7.49), txt_top=Inches(2.16), txt_width=Inches(0.65), txt_height=Inches(0.5),
        font_size=12, bold=True
    ),
    TopItemSpec(
        img_left=Inches(8.45), img_top=Inches(1.65), img_width=Inches(0.44), img_height=Inches(0.44),
        txt_left=Inches(8.48), txt_top=Inches(2.16), txt_width=Inches(0.65), txt_height=Inches(0.5),
        font_size=12, bold=True
    ),
]

# ---------- Helpers de dibujo ----------

def _add_text(slide, text: str, *, left, top, width, height, font_size, bold, font_name, font_color):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bool(bold)
    run.font.name = font_name
    run.font.color.rgb = RGBColor.from_string("2E4D58")  # Color hexadecimal
    


def _resolve_image_path(p: str | Path, img_dir: Optional[Path]) -> Path:
    p = Path(p)
    if p.is_absolute():
        return p
    if p.exists():
        return p
    if img_dir is None:
        return p
    return img_dir / p

# ---------- Preparación del top-5 desde longtable ----------

def build_topn_by_id(
    df_long: pd.DataFrame,
    *,
    id_col: str = "id",
    score_col: str = "score",
    #label_col: str = "label",
    image_col: str = "image",
    n: int = 5,
    ascending: bool = False,
) -> Dict[object, List[dict]]:
    """
    Devuelve dict: {id: [ {label, score, image}, ... ] } con top-n por id.
    """
    work = df_long[[id_col, score_col, image_col]].copy()
    top = work.groupby(id_col, sort=False).head(n)

    out: Dict[object, List[dict]] = {}
    for _id, g in top.groupby(id_col, sort=False):
        out[_id] = g.to_dict(orient="records")
    return out

# ---------- Render principal (una pasada por slides) ----------

def render_slides_for_ids(
    prs: Presentation,
    df_main: pd.DataFrame,
    df_long: pd.DataFrame,
    *,
    id_col: str = "id",
    start_slide: int = 0,
    sort_by_id: bool = False,          # orden de slides contra df_main
    fixed_text_fields: List[TextSpec] = FIXED_TEXT_FIELDS,
    top_slots: List[TopItemSpec] = TOP5_SLOTS,
    # columnas del longtable
    long_score_col: str = "score",
    long_image_col: str = "image",
    top_n: int = 5,
    top_ascending: bool = False,
    img_dir: Optional[Path] = None,     # si image en df_long es nombre relativo
) -> Presentation:
    
    if id_col not in df_main.columns:
        raise ValueError(f"df_main no tiene '{id_col}'. Tiene: {list(df_main.columns)}")

    # Validar columnas fijas
    for spec in fixed_text_fields:
        if spec.col not in df_main.columns:
            raise ValueError(f"df_main no tiene columna '{spec.col}' requerida por TextSpec.")

    work = df_main.copy()
    if sort_by_id:
        work = work.sort_values(by=id_col, kind="mergesort")

    needed = len(work)
    available = len(prs.slides) - start_slide
    if available < needed:
        raise IndexError(f"No hay suficientes slides: necesitas {needed} desde {start_slide}, hay {available}.")

    top_by_id = build_topn_by_id(
        df_long,
        id_col=id_col,
        score_col=long_score_col,
        image_col=long_image_col,
        n=top_n,
        ascending=top_ascending,
    )

    # Una pasada por slides
    for i, row in enumerate(work.to_dict(orient="records")):
        the_id = row[id_col]
        slide = prs.slides[start_slide + i]

        # ---- Textos fijos (df_main) ----
        for spec in fixed_text_fields:
            val = row.get(spec.col)
            text = "" if pd.isna(val) else spec.fmt.format(val=val, **row)
            _add_text(
                slide,
                text,
                left=spec.left, top=spec.top, width=spec.width, height=spec.height,
                font_size=spec.font_size, bold=spec.bold, font_name=spec.font_name,
                font_color=spec.font_color,
            )

        # ---- Top-N (df_long) ----
        items = top_by_id.get(the_id, [])
        # Relleno: si hay menos de 5, deja slots vacíos sin romper
        for j, slot in enumerate(top_slots):
            if j >= len(items):
                continue

            item = items[j]
            img_path = _resolve_image_path(item[long_image_col], img_dir)
            if not img_path.exists():
                # Decide tu política: aquí fallo fuerte para que lo detectes.
                raise FileNotFoundError(f"No existe imagen para id={the_id}, slot={j}: {img_path}")

            slide.shapes.add_picture(
                str(img_path),
                slot.img_left, slot.img_top,
                width=slot.img_width, height=slot.img_height
            )

            # texto del item (label + score, etc.)
            score = item[long_score_col]
            txt = slot.fmt.format( score=score, id=the_id, **row)

            _add_text(
                slide,
                txt,
                left=slot.txt_left, top=slot.txt_top, width=slot.txt_width, height=slot.txt_height,
                font_size=slot.font_size, bold=slot.bold, font_name=slot.font_name,font_color=spec.font_color,
            )

    return prs
    
    
## ---------------- funciones para carga ----------------###

from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request


# función para autenticarse 

SCOPES = ["https://www.googleapis.com/auth/drive.file"]
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

def get_drive_service_oauth(
    client_secret_json: str | Path,
    token_path: str | Path = "token_drive.json",
):
    client_secret_json = Path(client_secret_json).resolve()
    token_path = Path(token_path).resolve()

    creds = None
    if token_path.exists():
        creds = Credentials.from_authorized_user_file(str(token_path), SCOPES)

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            # Abre navegador la primera vez y guarda token
            flow = InstalledAppFlow.from_client_secrets_file(str(client_secret_json), SCOPES)
            creds = flow.run_local_server(port=0)
        token_path.write_text(creds.to_json(), encoding="utf-8")

    return build("drive", "v3", credentials=creds)

def upload_pptx_oauth(
    pptx_path: str | Path,
    *,
    folder_id: str,
    client_secret_json: str | Path,
    token_path: str | Path = "token_drive.json",
) -> dict:
    pptx_path = Path(pptx_path).resolve()
    if not pptx_path.exists():
        raise FileNotFoundError(pptx_path)

    service = get_drive_service_oauth(client_secret_json, token_path)

    file_metadata = {"name": pptx_path.name, "parents": [folder_id]}
    media = MediaFileUpload(str(pptx_path), mimetype=PPTX_MIME, resumable=True)

    req = service.files().create(
        body=file_metadata,
        media_body=media,
        fields="id,name,webViewLink",
        supportsAllDrives=True,
    )

    resp = None
    while resp is None:
        status, resp = req.next_chunk()

    return resp


